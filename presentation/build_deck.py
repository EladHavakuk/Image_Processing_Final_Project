"""
Builds the final presentation deck using python-pptx (pure Python, no Node.js
dependency). Mirrors the structure of the retired build_deck.js -- see README
troubleshooting log for why this project moved off pptxgenjs.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = Path(__file__).resolve().parent
ASSETS = HERE / "assets"

# ---- Palette ----
NAVY = RGBColor(0x1B, 0x2A, 0x41)
NAVY_LIGHT = RGBColor(0x2E, 0x42, 0x58)
OFFWHITE = RGBColor(0xF4, 0xF6, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
BLUE = RGBColor(0x2E, 0x86, 0xAB)
GREEN = RGBColor(0x5B, 0x8C, 0x5A)
TEXT_DARK = RGBColor(0x1B, 0x2A, 0x41)
TEXT_MUTED = RGBColor(0x5A, 0x6B, 0x7A)
ICE = RGBColor(0xCA, 0xDC, 0xFC)
BORDER = RGBColor(0xDD, 0xE3, 0xE9)
FOOTER_LIGHT = RGBColor(0xA0, 0xAD, 0xB8)
FOOTER_DARK = RGBColor(0x8A, 0x9B, 0xB0)
RED = RGBColor(0xC0, 0x39, 0x2B)

FONT_HEAD = "Cambria"
FONT_BODY = "Calibri"

PW, PH = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(PW)
prs.slide_height = Inches(PH)
BLANK = prs.slide_layouts[6]


def add_slide(bg=None):
    slide = prs.slides.add_slide(BLANK)
    if bg is not None:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg
    return slide


def no_line(shape):
    shape.line.fill.background()
    return shape


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=1.0, radius=None, shadow=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius is not None else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if radius is not None:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line_color is None:
        no_line(shp)
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    if shadow:
        shp.shadow.inherit = False
    else:
        shp.shadow.inherit = False
    return shp


def add_ellipse(slide, x, y, d, fill):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    no_line(shp)
    shp.shadow.inherit = False
    return shp


def set_text(tf, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0,
             font=FONT_BODY, wrap=True):
    """runs: list of (text, {size, bold, italic, color, font, strike}) OR a single string."""
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(runs, str):
        runs = [(runs, {})]
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for text, opts in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(opts.get("size", 14))
        r.font.bold = opts.get("bold", False)
        r.font.italic = opts.get("italic", False)
        r.font.name = opts.get("font", font)
        r.font.color.rgb = opts.get("color", TEXT_DARK)
        if opts.get("strike"):
            r.font._rPr.set("strike", "sngStrike")


def add_text(slide, x, y, w, h, runs, **kwargs):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(box.text_frame, runs, **kwargs)
    return box


def add_bullets(slide, x, y, w, h, items, color=TEXT_DARK, size=13.5, line_spacing=1.3, font=FONT_BODY):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(4)
        r = p.add_run()
        r.text = "\u25CF  " + item
        r.font.size = Pt(size)
        r.font.name = font
        r.font.color.rgb = color
    return box


def icon_circle(slide, x, y, d, fill_color, label, label_color=WHITE, font_size=20):
    add_ellipse(slide, x, y, d, fill_color)
    add_text(slide, x, y, d, d, [(label, {"size": font_size, "bold": True, "color": label_color, "font": FONT_HEAD})],
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def footer(slide, page_num, dark=False):
    add_text(slide, 0.5, PH - 0.45, PW - 1, 0.3,
              [(f"Vision Robustness Under Distortion  \u2022  {page_num}", {"size": 9, "color": FOOTER_DARK if dark else FOOTER_LIGHT})],
              align=PP_ALIGN.RIGHT)


def add_image(slide, path, x, y, w=None, h=None):
    slide.shapes.add_picture(str(path), Inches(x), Inches(y),
                              Inches(w) if w else None, Inches(h) if h else None)


# ============================================================
# Slide 1 - Title
# ============================================================
s = add_slide(NAVY)
add_ellipse(s, 10.6, -1.4, 4.5, NAVY_LIGHT)
add_ellipse(s, -1.6, 4.8, 3.6, NAVY_LIGHT)
add_text(s, 0.9, 2.35, 11.5, 1.0, [("ROBUSTNESS OF VISION ALGORITHMS", {"size": 40, "bold": True, "color": WHITE, "font": FONT_HEAD})])
add_text(s, 0.9, 3.15, 11.5, 0.7, [("Under Real-World Image Distortions", {"size": 26, "color": AMBER, "font": FONT_HEAD})])
add_text(s, 0.9, 4.05, 9.5, 1.0, [
    ("Classical & deep-learning vision tasks, evaluated on real driving imagery \u2014\n"
     "compression, low-light, and motion-blur distortions, classical restoration,\n"
     "and YOLO fine-tuning.", {"size": 15, "color": ICE})
], line_spacing=1.25)
add_text(s, 0.9, 6.55, 9, 0.4, [("Digital Image Processing & Computer Vision \u2014 Course Project", {"size": 12, "color": FOOTER_DARK})])

# ============================================================
# Slide 2 - Project at a glance
# ============================================================
s = add_slide(OFFWHITE)
add_text(s, 0.7, 0.45, 11.9, 0.7, [("Project at a Glance", {"size": 32, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
rows = [
    ("1", "Dataset", "BDD100K \u2014 150-image subset with real detection ground truth", BLUE),
    ("2", "Tasks", "Edge/Corner detection, Line detection (both low-level, classical) + Object detection (high-level, DL)", GREEN),
    ("3", "Distortions", "Compression, Low-light, Motion blur \u2014 all directly relevant to dashcam footage", AMBER),
    ("4", "Severity", "5 calibrated levels per distortion, measured in PSNR (dB)", BLUE),
    ("5", "Restoration", "Bilateral deblocking, Gamma+CLAHE, Wiener deconvolution (tuned per severity) \u2014 matched to each distortion's cause", GREEN),
    ("6", "Fine-tuning", "YOLOv8n only \u2014 the classical tasks have no trainable weights", AMBER),
]
y = 1.5
for num, label, desc, color in rows:
    icon_circle(s, 0.75, y, 0.55, color, num, font_size=20)
    add_text(s, 1.55, y - 0.06, 3.0, 0.5, [(label, {"size": 16, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
    add_text(s, 4.65, y - 0.06, 8.0, 0.7, [(desc, {"size": 13, "color": TEXT_MUTED})], line_spacing=1.15)
    y += 0.88
footer(s, 2)

# ============================================================
# Slide 3 - Pipeline architecture
# ============================================================
s = add_slide(OFFWHITE)
add_text(s, 0.7, 0.45, 11.9, 0.7, [("Four-Stage Pipeline", {"size": 32, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
add_text(s, 0.7, 1.05, 11, 0.4, [("Every image passes through all four stages, evaluated on all 3 tasks", {"size": 14, "color": TEXT_MUTED})])

stages = [
    ("1", "CLEAN", "Baseline evaluation\non the original image", BLUE),
    ("2", "DISTORTED", "3 distortions \u00d7 5 severity\nlevels applied", AMBER),
    ("3", "RESTORED", "Classical restoration\napplied to distorted image", GREEN),
    ("4", "FINE-TUNED", "YOLOv8n fine-tuned on\ndistorted images (detection only)", NAVY),
]
boxW, gap, startX, y0, boxH = 2.55, 0.45, 0.85, 2.6, 2.6
for i, (num, title, desc, color) in enumerate(stages):
    x = startX + i * (boxW + gap)
    add_rect(s, x, y0, boxW, boxH, fill=WHITE, line_color=BORDER, radius=0.08)
    icon_circle(s, x + boxW / 2 - 0.35, y0 + 0.3, 0.7, color, num, font_size=24)
    add_text(s, x + 0.1, y0 + 1.15, boxW - 0.2, 0.4, [(title, {"size": 15, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})], align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, y0 + 1.6, boxW - 0.3, 0.9, [(desc, {"size": 11.5, "color": TEXT_MUTED})], align=PP_ALIGN.CENTER, line_spacing=1.2)
    if i < len(stages) - 1:
        add_text(s, x + boxW + 0.05, y0 + 0.9, gap - 0.1, 0.6, [("\u2192", {"size": 26, "color": RGBColor(0xB7, 0xC2, 0xCC)})], align=PP_ALIGN.CENTER)

add_text(s, 0.85, 5.65, 11.6, 0.4, [
    ("100 images \u00d7 (1 clean + 3 distortions \u00d7 5 levels \u00d7 2 stages) \u2192 1,600 image variants across the pipeline",
     {"size": 12.5, "italic": True, "color": TEXT_MUTED})
], align=PP_ALIGN.CENTER)
footer(s, 3)

# ============================================================
# Slide 4 - Dataset
# ============================================================
s = add_slide(OFFWHITE)
add_text(s, 0.7, 0.45, 11.9, 0.7, [("Dataset: BDD100K", {"size": 32, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])

add_text(s, 0.7, 1.35, 5.6, 0.4, [("Why BDD100K", {"size": 16, "bold": True, "color": BLUE, "font": FONT_HEAD})])
add_bullets(s, 0.7, 1.8, 5.6, 1.8, [
    "Real driving footage with genuine detection ground truth",
    "Cityscapes and KITTI both required login-gated downloads that couldn't be scripted",
    "150-image subset, train split \u2014 small but sufficient for this scope",
], size=13.5, line_spacing=1.3)

add_text(s, 0.7, 3.75, 5.6, 0.4, [("Class distribution (150-image subset)", {"size": 16, "bold": True, "color": BLUE, "font": FONT_HEAD})])
classes = [("car", 1543), ("traffic light", 400), ("truck", 83), ("bus", 43), ("person", 13)]
max_v = 1543
cy = 4.25
for name, count in classes:
    bar_w = max(3.6 * (count / max_v), 0.08)
    add_text(s, 0.7, cy, 1.3, 0.32, [(name, {"size": 11.5, "color": TEXT_DARK})], anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s, 2.05, cy + 0.03, bar_w, 0.26, fill=AMBER, radius=0.3)
    add_text(s, 2.05 + bar_w + 0.1, cy, 0.9, 0.32, [(str(count), {"size": 11, "color": TEXT_MUTED})], anchor=MSO_ANCHOR.MIDDLE)
    cy += 0.42

add_image(s, ASSETS / "detection_overlay.png", 6.65, 1.4, w=6.0, h=2.36)
add_text(s, 6.65, 3.8, 6.0, 0.55, [
    ("GT sanity check \u2014 predictions (green) vs. ground truth (red), clean vs. severe low-light",
     {"size": 10.5, "italic": True, "color": TEXT_MUTED})
])

add_rect(s, 6.65, 4.55, 6.0, 1.75, fill=WHITE, line_color=BORDER, radius=0.06)
add_text(s, 6.95, 4.72, 5.4, 0.35, [("Known limitation", {"size": 13, "bold": True, "color": NAVY, "font": FONT_HEAD})])
add_text(s, 6.95, 5.1, 5.4, 1.1, [
    ("Zero bicycle/motorcycle/train instances in this sample; heavily skewed toward cars and traffic lights. "
     "Per-class metrics for rare classes aren't meaningful at this scale.", {"size": 12, "color": TEXT_MUTED})
], line_spacing=1.25)
footer(s, 4)

# ============================================================
# Slide 5 - Three tasks
# ============================================================
s = add_slide(OFFWHITE)
add_text(s, 0.7, 0.45, 11.9, 0.7, [("Three Tasks: Low-Level + High-Level", {"size": 32, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
tasks = [
    ("Edge / Corner\nDetection", "LOW-LEVEL \u00b7 CLASSICAL", "cv2.Canny + goodFeaturesToTrack (Shi-Tomasi) + ORB",
     "No GT exists \u2014 measured against the clean image's own output (edge IoU, corner ratio, ORB match ratio)", BLUE),
    ("Line\nDetection", "LOW-LEVEL \u00b7 CLASSICAL", "cv2.Canny + HoughLinesP",
     "Replaced an original optical-flow/tracking plan that needed video frames we didn't have", GREEN),
    ("Object\nDetection", "HIGH-LEVEL \u00b7 DEEP LEARNING", "ultralytics YOLOv8n (COCO-pretrained)",
     "The one task with real ground truth \u2014 scored with proper IoU-matched precision/recall/F1", AMBER),
]
cardW, gapc, startXc, y0, cardH = 3.75, 0.35, 0.7, 1.5, 5.25
for i, (title, tag, method, note, color) in enumerate(tasks):
    x = startXc + i * (cardW + gapc)
    add_rect(s, x, y0, cardW, cardH, fill=WHITE, line_color=BORDER, radius=0.06)
    icon_circle(s, x + 0.35, y0 + 0.35, 0.65, color, str(i + 1), font_size=22)
    add_text(s, x + 0.3, y0 + 1.15, cardW - 0.6, 0.75, [(title, {"size": 19, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})], line_spacing=1.05)
    add_rect(s, x + 0.3, y0 + 2.05, cardW - 0.6, 0.36, fill=color, radius=0.5)
    add_text(s, x + 0.3, y0 + 2.05, cardW - 0.6, 0.36, [(tag, {"size": 10.5, "bold": True, "color": WHITE})], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s, x + 0.3, y0 + 2.65, cardW - 0.6, 0.85, [(method, {"size": 12.5, "bold": True, "color": TEXT_DARK})], line_spacing=1.2)
    add_text(s, x + 0.3, y0 + 3.55, cardW - 0.6, 1.55, [(note, {"size": 11.5, "italic": True, "color": TEXT_MUTED})], line_spacing=1.25)
footer(s, 5)

# ============================================================
# Slide 6 - Distortions
# ============================================================
s = add_slide(OFFWHITE)
add_text(s, 0.7, 0.4, 11.9, 0.65, [("Three Distortions", {"size": 32, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
add_text(s, 0.7, 0.98, 11.9, 0.4, [
    ("Chosen for direct relevance to dashcam/driving conditions \u2014 5 calibrated severity levels each, measured in PSNR (dB)",
     {"size": 13, "color": TEXT_MUTED})
])
add_image(s, ASSETS / "before_after_grid.png", 3.55, 1.5, w=6.25, h=4.33)

legend = [
    ("Compression", "quality 50\u21923 \u2192 ~39\u219224 dB", BLUE),
    ("Low-light", "brightness -0.2\u2192-0.9 \u2192 ~18\u21927 dB", AMBER),
    ("Motion blur", "kernel 5\u219231px \u2192 ~33\u219225 dB", GREEN),
]
ly = 1.7
for name, rng, color in legend:
    add_rect(s, 0.7, ly, 0.28, 0.28, fill=color, radius=0.2)
    add_text(s, 1.12, ly - 0.06, 2.2, 0.4, [(name, {"size": 13.5, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
    add_text(s, 0.7, ly + 0.32, 2.55, 0.55, [(rng, {"size": 10.5, "color": TEXT_MUTED})], line_spacing=1.15)
    ly += 1.15
add_text(s, 0.55, 5.35, 2.7, 1.3, [
    ("Motion blur uses a known, controlled linear kernel (not random) \u2014 so restoration can be genuine non-blind deconvolution.",
     {"size": 10.5, "italic": True, "color": TEXT_MUTED})
], line_spacing=1.25)
footer(s, 6)

# ============================================================
# Slide 7 - Restoration
# ============================================================
s = add_slide(OFFWHITE)
add_text(s, 0.7, 0.45, 11.9, 0.7, [("Restoration Methods (Stage 3)", {"size": 32, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
methods = [
    ("Compression", "Bilateral filter on Y (luma) channel only", "Smooths blocking artifacts while preserving color and most edges", BLUE),
    ("Low-light", "Gamma correction (\u03b3=0.35) + CLAHE on L channel", "Lifts shadow detail, then boosts local contrast without blowing out highlights", AMBER),
    ("Motion blur", "Wiener deconvolution, regularization tuned per severity level", "Best detection F1 of 3 methods explicitly compared (0.357 vs. 0.323 for no restoration) \u2014 see next slide", GREEN),
]
y = 1.6
for name, method, why, color in methods:
    add_rect(s, 0.7, y, 11.9, 1.5, fill=WHITE, line_color=BORDER, radius=0.06)
    icon_circle(s, 1.0, y + 0.45, 0.6, color, name[0], font_size=20)
    add_text(s, 1.85, y + 0.18, 2.6, 0.5, [(name, {"size": 16, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
    add_text(s, 1.85, y + 0.68, 2.6, 0.7, [(method, {"size": 10.5, "bold": True, "color": color})], line_spacing=1.15)
    add_text(s, 4.7, y + 0.28, 7.6, 0.95, [(why, {"size": 12.5, "color": TEXT_MUTED})], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.25)
    y += 1.7
footer(s, 7)

# ============================================================
# Slide 8 - Restoration method comparison (motion blur ablation)
# ============================================================
s = add_slide(WHITE)
add_text(s, 0.7, 0.4, 11.9, 0.65, [("Restoration Method Comparison: Motion Blur", {"size": 28, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
add_text(s, 0.7, 0.95, 11.9, 0.4, [
    ("Three methods, kept side by side on purpose \u2014 no single metric predicted the best one for the actual task", {"size": 13, "color": TEXT_MUTED})
])
add_image(s, ASSETS / "motion_blur_method_grid.png", 0.5, 1.4, w=12.3, h=3.55)

summary_rows = [
    ("Method", "Det. F1", "PSNR", "Stripe", ""),
    ("Distorted (none)", "0.323", "27.66", "0.54", ""),
    ("Wiener (fixed)", "0.311", "25.70", "2.78", "collapses at severe level (F1=0.017)"),
    ("Wiener (tuned)", "0.357", "28.31", "1.49", "production choice \u2014 best detection F1"),
    ("Richardson-Lucy", "0.291", "27.58", "2.04", "cleanest-looking, worst for detection"),
]
colx = [0.7, 3.6, 5.0, 6.3, 7.5]
colw = [2.9, 1.4, 1.3, 1.2, 5.0]
ty = 5.15
for r, row in enumerate(summary_rows):
    is_header = r == 0
    for c, (text, x, w) in enumerate(zip(row, colx, colw)):
        color = TEXT_MUTED if is_header else TEXT_DARK
        bold = is_header
        size = 11 if is_header else 11.5
        add_text(s, x, ty, w, 0.32, [(text, {"size": size, "bold": bold, "color": color, "italic": (c == 4 and not is_header)})])
    ty += 0.34
add_rect(s, 0.65, 5.1, 11.95, 0.02, fill=BORDER)
footer(s, 8)

# ============================================================
# Slide 9 - Results: Object Detection
# ============================================================
s = add_slide(WHITE)
add_text(s, 0.7, 0.4, 11.9, 0.65, [("Results: Object Detection Robustness", {"size": 30, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
add_text(s, 0.7, 0.98, 11.9, 0.4, [
    ("Detection F1 vs. real BDD100K ground truth \u2014 solid = distorted, dashed = after restoration", {"size": 13, "color": TEXT_MUTED})
])
add_image(s, ASSETS / "det_f1_slide.png", 1.55, 1.55, w=10.2, h=5.58)
footer(s, 9)

# ============================================================
# Slide 9 - Results: low-level tasks
# ============================================================
s = add_slide(WHITE)
add_text(s, 0.7, 0.4, 11.9, 0.65, [("Results: Edge/Corner & Line Detection", {"size": 30, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
add_text(s, 0.7, 0.98, 11.9, 0.4, [
    ("Measured against each image's own clean-image output (no external GT available for these low-level tasks)", {"size": 13, "color": TEXT_MUTED})
])
add_image(s, ASSETS / "edge_iou_slide.png", 0.35, 1.6, w=6.3, h=3.45)
add_image(s, ASSETS / "line_iou_slide.png", 6.7, 1.6, w=6.3, h=3.45)
add_rect(s, 0.7, 5.3, 11.9, 1.4, fill=OFFWHITE, radius=0.06)
add_text(s, 1.0, 5.45, 11.3, 1.1, [
    ("Notable: ", {"bold": True, "color": NAVY, "size": 13.5}),
    ("with the production restoration method (tuned Wiener), motion blur improves on both fronts at every severity level "
     "\u2014 edge IoU 0.44\u21920.51 and detection F1 0.32\u21920.36 \u2014 but that wasn't true for every method tried; see the dedicated comparison slide.",
     {"color": TEXT_MUTED, "size": 13.5}),
], anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.3)
footer(s, 10)

# ============================================================
# Slide 10 - Key findings
# ============================================================
s = add_slide(NAVY)
add_text(s, 0.7, 0.45, 11.9, 0.7, [("Key Findings", {"size": 32, "bold": True, "color": WHITE, "font": FONT_HEAD})])
findings = [
    ("Clear degradation with severity", "Compression detection F1 falls from 0.45 (mild) to 0.15 (severe) as SNR drops.", AMBER),
    ("Restoration helps most when things are worst", "Deblocking barely changes mild-compression F1, but meaningfully helps at severe levels.", BLUE),
    ("Pixel quality \u2260 downstream performance", "Low-light restoration improves SNR (9.4\u219212.2 dB) but slightly hurts detection F1 \u2014 independently re-verified, next slide.", GREEN),
    ("\u201cCleanest-looking\u201d restoration isn't always best", "Richardson-Lucy looks cleanest but has the worst detection F1 of 3 methods (0.291); tuned Wiener wins on the metric that matters (0.357) \u2014 also re-verified, next slide.", AMBER),
    ("Baseline domain gap motivates fine-tuning", "Even on clean images, stock YOLOv8n only reaches F1=0.45 against BDD100K's real GT.", BLUE),
]
y = 1.4
for i, (title, desc, color) in enumerate(findings):
    icon_circle(s, 0.75, y, 0.5, color, str(i + 1), label_color=NAVY, font_size=17)
    add_text(s, 1.5, y - 0.08, 10.8, 0.4, [(title, {"size": 16, "bold": True, "color": WHITE, "font": FONT_HEAD})])
    add_text(s, 1.5, y + 0.33, 10.8, 0.55, [(desc, {"size": 12.5, "color": ICE})], line_spacing=1.2)
    y += 1.08
footer(s, 11, dark=True)

# ============================================================
# Slide 11b - Verification: why restoration can score worse
# ============================================================
s = add_slide(WHITE)
add_text(s, 0.7, 0.4, 11.9, 0.65, [("Verification: Why Restoration Can Score Worse", {"size": 27, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
add_text(s, 0.7, 0.98, 11.9, 0.45, [
    ("\u201cRestoration hurts detection\u201d sounds like a bug by default \u2014 checked directly rather than assumed", {"size": 13, "color": TEXT_MUTED})
])

add_rect(s, 0.7, 1.6, 5.75, 0.9, fill=OFFWHITE, radius=0.08)
add_text(s, 0.95, 1.75, 5.3, 0.6, [
    ("Ruled out: data corruption", {"bold": True, "size": 13, "color": NAVY}),
    ("\nEvery restored image keeps the exact same shape/dtype as the input \u2014 confirmed programmatically, not assumed.", {"size": 11, "color": TEXT_MUTED}),
], line_spacing=1.2)

add_rect(s, 6.85, 1.6, 5.75, 0.9, fill=OFFWHITE, radius=0.08)
add_text(s, 7.1, 1.75, 5.3, 0.6, [
    ("The mechanism: verified with Laplacian variance", {"bold": True, "size": 13, "color": NAVY}),
    ("\nA standard noise/detail proxy \u2014 real structure AND injected noise both raise it.", {"size": 11, "color": TEXT_MUTED}),
], line_spacing=1.2)

cardY = 2.7
add_rect(s, 0.7, cardY, 5.75, 3.05, fill=WHITE, line_color=BORDER, radius=0.08)
add_text(s, 0.95, cardY + 0.18, 5.3, 0.35, [("Low-light: CLAHE", {"bold": True, "size": 15, "color": BLUE, "font": FONT_HEAD})])
add_text(s, 0.95, cardY + 0.6, 5.3, 2.3, [
    ("2 correct detections \u2192 ", {"size": 12, "color": TEXT_DARK}),
    ("0", {"size": 12, "bold": True, "color": RED}),
    (" after restoration, despite brightness genuinely improving (11.3\u219238.0).\n\n", {"size": 12, "color": TEXT_DARK}),
    ("Laplacian variance in that image: ", {"size": 12, "color": TEXT_DARK}),
    ("clean=456, distorted=211, restored=1451", {"size": 12, "bold": True, "color": BLUE}),
    (" \u2014 CLAHE injects 3\u00d7 more high-frequency noise than even the clean image has.\n\n", {"size": 12, "color": TEXT_DARK}),
    ("Checked at scale (40 images): restored texture exceeds clean's own level in ", {"size": 12, "color": TEXT_DARK}),
    ("35/40 images (87.5%)", {"size": 12, "bold": True, "color": BLUE}),
    (".", {"size": 12, "color": TEXT_DARK}),
], line_spacing=1.28)

add_rect(s, 6.85, cardY, 5.75, 3.05, fill=WHITE, line_color=BORDER, radius=0.08)
add_text(s, 7.1, cardY + 0.18, 5.3, 0.35, [("Motion blur: Richardson-Lucy", {"bold": True, "size": 15, "color": GREEN, "font": FONT_HEAD})])
add_text(s, 7.1, cardY + 0.6, 5.3, 2.3, [
    ("7 detections (distorted) \u2192 5, ", {"size": 12, "color": TEXT_DARK}),
    ("higher confidence", {"size": 12, "bold": True, "color": GREEN}),
    (" (tuned Wiener) vs. \u2192 ", {"size": 12, "color": TEXT_DARK}),
    ("1, low confidence", {"size": 12, "bold": True, "color": RED}),
    (" (Richardson-Lucy).\n\n", {"size": 12, "color": TEXT_DARK}),
    ("Laplacian variance: distorted=48, wiener_tuned=13 (smooths), richardson_lucy=57", {"size": 12, "bold": True, "color": GREEN}),
    (" \u2014 RL adds MORE high-frequency noise than the original blur.\n\n", {"size": 12, "color": TEXT_DARK}),
    ("Checked at scale (40 images): RL exceeds the distorted image's own noise level in ", {"size": 12, "color": TEXT_DARK}),
    ("39/39 images (100%)", {"size": 12, "bold": True, "color": GREEN}),
    (".", {"size": 12, "color": TEXT_DARK}),
], line_spacing=1.28)

add_text(s, 0.7, 5.95, 11.9, 0.9, [
    ("Honest caveat: ", {"bold": True, "size": 11.5, "color": NAVY}),
    ("noise magnitude only weakly predicts F1 drop per-image (r \u2248 \u22120.13). The mechanism is real and consistent in aggregate \u2014 restoration methods optimize for pixel fidelity, not for what a frozen pretrained detector finds recognizable \u2014 but a detector's response to it is nonlinear, not a simple proportional rule.", {"size": 11.5, "color": TEXT_MUTED}),
], line_spacing=1.25)
footer(s, 12)

# ============================================================
# Slide 11 - Fine-tuning
# ============================================================
s = add_slide(WHITE)
add_text(s, 0.7, 0.4, 11.9, 0.65, [("Fine-Tuning (Stage 4)", {"size": 30, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
add_image(s, ASSETS / "finetune_slide.png", 0.6, 1.15, w=6.6, h=3.81)

add_rect(s, 7.55, 1.2, 5.1, 1.75, fill=OFFWHITE, radius=0.06)
add_text(s, 7.85, 1.35, 4.5, 0.3, [("Overall Detection F1", {"size": 12, "color": TEXT_MUTED})])
add_text(s, 7.85, 1.68, 4.5, 0.65, [
    ("0.267 ", {"size": 26, "color": TEXT_MUTED, "strike": True, "font": FONT_HEAD}),
    (" \u2192  ", {"size": 22, "color": TEXT_MUTED, "font": FONT_HEAD}),
    ("0.261", {"size": 30, "bold": True, "color": AMBER, "font": FONT_HEAD}),
], anchor=MSO_ANCHOR.MIDDLE)
add_text(s, 7.85, 2.4, 4.5, 0.4, [("baseline \u2192 fine-tuned (held-out distorted set)", {"size": 11, "italic": True, "color": TEXT_MUTED})])

add_rect(s, 7.55, 3.15, 5.1, 2.45, fill=WHITE, line_color=BORDER, radius=0.06)
add_text(s, 7.85, 3.35, 4.5, 0.4, [("Honest result: no clear win", {"size": 14.5, "bold": True, "color": NAVY, "font": FONT_HEAD})])
add_text(s, 7.85, 3.78, 4.5, 1.75, [
    ("F1 is essentially flat overall \u2014 small recall gain traded for a precision drop. Results are mixed per distortion "
     "(compression improved slightly, low-light got worse).\n\nExpected at this scale: 40\u201350 training images, 5\u20138 epochs, "
     "CPU-only \u2014 squarely the regime where catastrophic forgetting can outweigh adaptation.",
     {"size": 11.5, "color": TEXT_MUTED})
], line_spacing=1.25)
footer(s, 13)

# ============================================================
# Slide 12 - Process highlights
# ============================================================
s = add_slide(OFFWHITE)
add_text(s, 0.7, 0.45, 11.9, 0.6, [("Process Highlights", {"size": 32, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})])
add_text(s, 0.7, 1.02, 11.9, 0.4, [("Real friction points hit while building this \u2014 documented, not hidden", {"size": 13, "color": TEXT_MUTED})])
items = [
    ("Dataset access", "Cityscapes and KITTI both required login-gated downloads. Switched to BDD100K, downloaded manually, filtered locally.", BLUE),
    ("1GB label file", "BDD100K's full label file (~70k entries) was too big to hand off. A local script filtered it down to the 150 images actually used, in under 1MB.", GREEN),
    ("The fine-tuning caching bug", "First attempt showed byte-identical baseline vs. fine-tuned results \u2014 a cached YOLO model got mutated in-place by .train(). Fixed with two independent model instances.", AMBER),
    ("A restoration bug, chased across 3 attempts", "Fixed-balance Wiener caused striping; tuning it per-level fixed that. Richardson-Lucy then looked even better on one test image \u2014 but a full 150-image comparison showed it actually has the worst detection F1 of all 3. Tuned Wiener (the one abandoned in step 2) is the real winner.", BLUE),
]
y = 1.5
for title, desc, color in items:
    add_rect(s, 0.7, y, 11.9, 1.2, fill=WHITE, line_color=BORDER, radius=0.06)
    icon_circle(s, 0.95, y + 0.3, 0.5, color, "!", font_size=19)
    add_text(s, 1.7, y + 0.1, 2.85, 1.0, [(title, {"size": 13.5, "bold": True, "color": TEXT_DARK, "font": FONT_HEAD})], line_spacing=1.1)
    add_text(s, 4.75, y + 0.1, 7.55, 1.0, [(desc, {"size": 10.5, "color": TEXT_MUTED})], line_spacing=1.14)
    y += 1.35
footer(s, 14)

# ============================================================
# Slide 13 - Limitations & closing
# ============================================================
s = add_slide(NAVY)
add_ellipse(s, 10.9, 4.9, 4.0, NAVY_LIGHT)
add_text(s, 0.9, 0.6, 11.5, 0.65, [("Limitations & Takeaways", {"size": 30, "bold": True, "color": WHITE, "font": FONT_HEAD})])
lims = [
    "Small scale by design (150 images) \u2014 explicitly permitted for this course project; pipeline scales to any N",
    "Class imbalance: zero bicycle/motorcycle/train instances in this sample",
    "Edge/corner and line detection use clean-image-as-reference, not external ground truth",
    "BDD100K \u2192 COCO class mapping is approximate (\u201crider\u201d folds into \u201cperson\u201d; \u201ctraffic sign\u201d dropped)",
    "Fine-tuning is a proof of concept, not a production training run",
]
add_bullets(s, 0.9, 1.5, 10.5, 2.6, lims, color=ICE, size=14, line_spacing=1.4)

add_rect(s, 0.9, 4.5, 8.5, 0.02, fill=NAVY_LIGHT)
add_text(s, 0.9, 4.85, 8, 0.4, [("Full code, data, results, and this report:", {"size": 14, "color": FOOTER_DARK})])
add_text(s, 0.9, 5.3, 8, 0.55, [("github.com  \u2014  see README.md", {"size": 22, "bold": True, "color": AMBER, "font": FONT_HEAD})])
add_text(s, 0.9, 6.35, 8, 0.5, [("Thank you", {"size": 16, "italic": True, "color": WHITE, "font": FONT_HEAD})])

out_path = HERE / "vision_robustness_presentation.pptx"
prs.save(str(out_path))
print("Saved", out_path)
